# -*- coding: utf-8 -*-
"""
AI-PM-Curriculum PPT 生成器（数据驱动骨架）
============================================
设计目标：先搭骨架，方便往里塞。
- 改内容：只改下面 CHAPTERS / PANORAMA / APPENDIX / META 这几个数据区，不用动渲染逻辑。
- 加一页：往 CHAPTERS 列表里加一个 dict 即可。
- 换配色：改 COLORS 字典。
- 输出：<repo>/ppt/AI-PM-Curriculum.pptx

运行：python ppt/build_ppt.py
"""
import os
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─────────────────────────── 数据区（往这里塞内容）───────────────────────────
REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))  # 仓库根

META = {
    "title": "产品经理就业知识体系",
    "subtitle": "AI 增补版 · 课程全景",
    "footer": "产品经理就业知识体系 · AI 增补版",
    "version": "v0.2.0 · 2026-08",
    "out": os.path.join(REPO, "ppt", "AI-PM-Curriculum.pptx"),
}

# 全景：两列
PANORAMA = {
    "base": {
        "title": "通用底座  product/",
        "tag": "人人都会",
        "items": [
            ("①", "需求", "翻译"),
            ("②", "规划设计", "落地"),
            ("③", "项目管理", "交付"),
            ("④", "数据分析", "验证"),
            ("⑤", "商业战略", "方向"),
            ("⑥", "沟通协作", "放大器"),
            ("⑦", "学习与自我修养", "底色"),
        ],
    },
    "ai": {
        "title": "AI 增补层  skill/",
        "tag": "AI 岗位增量",
        "items": [
            ("⑧", "AI 技术认知", "能不能做"),
            ("⑨", "模型评估与数据闭环", "做得好不好"),
            ("⑩", "信任安全与伦理合规", "守住底线"),
            ("⑪", "AI 产品经济学", "值不值 / 赚不赚"),
        ],
    },
}

# 11 个板块，每页一个。layer: base=蓝 / ai=青
CHAPTERS = [
    dict(num="①", title="需求能力", focus="翻译", layer="base",
         summary="需求能力 = 翻译能力，从现象到本质，从「用户要的」到「真正该做的」。",
         terms=["用户需求 vs 产品需求：用户要解决方案，产品要真实需求",
                "需求采集四象限：定性 / 定量 × 说 / 做",
                "DNA 检测：定属性 → 商业价值 → 实现难度 → 性价比",
                "性价比 = 商业价值 ÷ 实现难度（开发人天）",
                "分·总·分：拆属性 → 归共性 → 回场景验证"],
         note=("AI 增量", "板块⑧把「翻译」升级为「翻译成模型任务」——多一个判断维度（能不能用 AI 做）。"),
         img="product/01-需求能力/思维导图.png"),
    dict(num="②", title="规划设计", focus="落地", layer="base",
         summary="规划设计 = 落地能力，把抽象变具体，把想法变可交付的约定。",
         terms=["产品 / 产品经理：发现并描述问题，转化为可交付的产品",
                "大产品观：凡影响体验与公司价值的一切都属 PM 职责",
                "BRD · MRD · PRD · FSD：值不值 → 给谁 → 做什么 → 怎么做",
                "PRD 七大板块：目标 / 范围 / 场景 / 功能 / 非功能 / 验收 / 遗留",
                "文档只是手段，服务于沟通与传承，不为写而写"],
         note=("改造点", "AI-native PRD 五要素：评估标准 / 模型约束 / 数据需求 / 埋点清单 / 防护兜底——核心多问一句「模型出错了怎么办」。"),
         img="product/02-规划设计/思维导图.png"),
    dict(num="③", title="项目管理", focus="交付", layer="base",
         summary="项目管理 = 交付能力，多快好省的平衡术，把承诺变成结果。",
         terms=["项目：只做一次、包含多项互相关联任务的工作",
                "多快好省（TRQ）：范围 / 时间 / 品质 / 资源不可兼得",
                "WBS：任务自上而下分解到可分配可估算的粒度",
                "三次评审：需求 / 设计 / 测试，防病优于治病",
                "敏捷开发：迭代范围不变、小步快跑、快速反馈"],
         note=("改造点", "AI 交付与运营：分层发布（alpha→beta→GA）/ 回滚降级 / Kill Switch / Prompt 版本控制 / 漂移监控——AI 上线只是开始。"),
         img="product/03-项目管理/思维导图.png"),
    dict(num="④", title="数据分析", focus="验证", layer="base",
         summary="数据分析 = 验证能力，从拍脑袋决策到用数据说话。",
         terms=["KPI：战略的具体表现，本质是手段而非目的",
                "北极星指标：一个能代表产品长期价值的核心指标",
                "漏斗与转化：每层流失都是优化机会，先查流失最大的一层",
                "留存 / 活跃：留存是产品价值的照妖镜",
                "A/B 测试：随机分两组对照，一次只改一个变量"],
         note=("改造点", "模型指标看板：离线指标（F1 / 准确率）+ 在线指标（采纳率 / 任务完成率 / 人工接管率），两套都要盯。"),
         img="product/04-数据分析/思维导图.png"),
    dict(num="⑤", title="商业战略", focus="方向", layer="base",
         summary="商业战略 = 方向能力，选对战场比打赢一场战斗更重要。",
         terms=["可行性三步曲：我们在哪儿 → 去哪儿 → 怎么去",
                "PEST：政治 / 经济 / 社会 / 技术四维扫外部环境",
                "SWOT：优势 / 劣势 / 机会 / 威胁，落到策略组合",
                "价值观 → 使命 → 愿景 → 战略：层层炼成，先于执行",
                "炮灰版 / 水平营销：价格锚点 + 新维度"],
         note=("改造点", "数据飞轮与防御性（详见板块⑪）：AI 产品最值钱的资产是数据飞轮，不是算法本身。"),
         img="product/05-商业战略/思维导图.png"),
    dict(num="⑥", title="沟通协作", focus="影响力", layer="base",
         summary="沟通协作 = 影响力能力，靠专业、靠谱、让人愿意跟你干。",
         terms=["无授权领导：不靠职位权力，靠专业话语权 + 影响力",
                "矩阵型组织：职能 + 项目双线融合，可能有双头领导",
                "接口人：团队间单一对接人，过滤噪音",
                "管理靠权力，领导靠魅力",
                "猎人 vs 农民：猎人管事盯结果，农民管人耕过程"],
         note=("改造点", "AI 沟通：对技术讲业务目标与验收口径，对高管讲单位经济，对销售讲使用边界与兜底，对合规讲风险前置。"),
         img="product/06-沟通协作/思维导图.png"),
    dict(num="⑦", title="学习与自我修养", focus="底色", layer="base",
         summary="自我修养 = 底色能力，技能可以速成，底色决定天花板。",
         terms=["自我修养四件套：爱生活 / 有理想 / 会思考 / 能沟通",
                "解决问题通用思路：万能提问模板，任何问题都能套",
                "少做就是多做：用 100% 质量做 75% 数量",
                "生态隐喻：云雨 → 河流 → 动植物 → 阳光 → 大地",
                "产品经理主义：把产品思维抽象成可普适的做事方法"],
         note=("改造点", "持续追踪 AI 进展：跟模型发布节奏走、订阅信源、动手做小实验——「上个月做不到的，这周可能就能做了」。"),
         img="product/07-学习与自我修养/思维导图.png"),
    dict(num="⑧", title="AI 技术认知", focus="判断能不能做", layer="ai",
         summary="AI 技术认知 = 判断能力，把「我有个 AI 想法」变成「技术上成立、边界在哪」。",
         terms=["模型类型与能力边界：边界内做深做透，边界外别硬做",
                "大模型基础：Transformer / 预训练 / 微调 / Prompt 工程",
                "RAG：先检索再生成，把「背题」变「开卷考试」",
                "Agent / MCP / 工具调用：模型自主规划并调用工具",
                "幻觉：一本正经地编造，无法根除，只能缓解"],
         note=("面试怎么考", "给一个 AI 产品 idea 问「能不能做」→ 翻译成模型任务 → 对能力边界 → 判断可行性 + 列风险点（幻觉 / 延迟 / 成本）。"),
         img="skill/08-AI技术认知/思维导图.png"),
    dict(num="⑨", title="模型评估与数据闭环", focus="验证做得好不好", layer="ai",
         summary="模型评估与数据闭环 = 验证能力，从「感觉还行」到「有数字证明」。",
         terms=["评估集 / 金标准集：人工标注的标准答案，不参与训练",
                "离线指标：准确率 / 召回率 / F1，评估集上跑",
                "在线指标：采纳率 / 任务完成率 / 人工接管率",
                "数据闭环：采集 → 清洗 → 标注 → 评估 → 迭代",
                "价值归因：判断指标变化是模型还是别的因素"],
         note=("面试怎么考", "问「AI 功能怎么验收」→ 先定评估集和离线指标 → 上线看在线指标 → 数据闭环迭代 → 归因到模型。"),
         img="skill/09-模型评估与数据闭环/思维导图.png"),
    dict(num="⑩", title="信任安全与伦理合规", focus="守住底线", layer="ai",
         summary="信任安全与伦理合规 = 底线能力，先想清楚怎么不闯祸，再谈功能多强。",
         terms=["公平性：对不同人群的表现不能有明显偏差",
                "可解释性：为什么给这个结果，能说清楚",
                "隐私保护：最小化采集 / 脱敏 / 授权 / 合规前置",
                "内容安全：前置过滤 + 后置拦截 + 人工巡查",
                "Red Teaming + 发布门槛：主动找漏洞，不过关不发"],
         note=("面试怎么考", "问「AI 上线前怎么把关」→ 红队 + 安全检查 + 兜底 + 发布门槛 + 分层发布。"),
         img="skill/10-信任安全与伦理合规/思维导图.png"),
    dict(num="⑪", title="AI 产品经济学", focus="算清值不值 / 赚不赚", layer="ai",
         summary="AI 产品经济学 = 算账能力，从「觉得有用」到「算得清赚不赚钱」。",
         terms=["Token / 推理成本：每次请求都产生真实成本",
                "单位经济模型：收入 − 成本，LTV > CAC + 累计推理成本",
                "定价建模：成本打底 + 价值锚定 + 版本分层",
                "模型选型权衡：成本 / 延迟 / 精度 / 上下文四维",
                "数据飞轮与防御性：越用数据越多 → 模型越好 → 用户越多"],
         note=("面试怎么考", "问「这个 AI 功能值不值得做」→ 单位经济 + 成本结构 + 飞轮潜力，三件套。"),
         img="skill/11-AI产品经济学/思维导图.png"),
]

APPENDIX = [
    ("A", "能力自评清单", "11 个板块各一栏自评（1–5 分）+ 补强动作，用于定期查漏。"),
    ("B", "面试高频问题快答", "通用底座 10 问（产品是什么 / 需求怎么排 / BRD-PRD 区别 / 项目延期 / 无职权推动…）+ AI 增补层 6 问（能不能做 / 准不准 / 答错怎么办 / 怎么定价 / 护城河 / 上线把关）。"),
    ("C", "优化路线（React 化）", "V0 想法定稿 → V1 内容数字化 → V2 React 应用 → V3 背得下来（记忆卡片 / 自测）→ V4 持续优化闭环。把知识体系当产品运营，issues 就是用户反馈。"),
]

# ─────────────────────────── 配色 / 版式常量 ───────────────────────────
COLORS = {
    "navy":  RGBColor(0x1B, 0x2A, 0x4A),
    "blue":  RGBColor(0x2D, 0x6C, 0xDF),
    "teal":  RGBColor(0x0F, 0xA3, 0xA3),
    "ink":   RGBColor(0x1F, 0x29, 0x37),
    "muted": RGBColor(0x6B, 0x72, 0x80),
    "bg":    RGBColor(0xF4, 0xF6, 0xFA),
    "card":  RGBColor(0xFF, 0xFF, 0xFF),
    "line":  RGBColor(0xE5, 0xE7, 0xEB),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}
LAYER_TINT = {"base": RGBColor(0xE8, 0xF0, 0xFE), "ai": RGBColor(0xE0, 0xF5, 0xF5)}
FONT = "Microsoft YaHei"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


def _set_font(run, size, bold=False, color=None, name=FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    # 东亚洲字体（中文回退）
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)


def _solid(shape, fill, line=None, line_w=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w or 0.75)
    shape.shadow.inherit = False


def _rect(slide, L, T, W, H, fill, line=None, round_=False, line_w=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(kind, Inches(L), Inches(T), Inches(W), Inches(H))
    if round_:
        try:
            sp.adjustments[0] = 0.08
        except Exception:
            pass
    _solid(sp, fill, line, line_w)
    return sp


def _text(slide, L, T, W, H, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """lines: list of paragraphs；每个 paragraph = list of (text, size, bold, color)"""
    tb = slide.shapes.add_textbox(Inches(L), Inches(T), Inches(W), Inches(H))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.0
        for (t, size, bold, color) in para:
            r = p.add_run()
            r.text = t
            _set_font(r, size, bold, color)
    return tb


def _footer(slide, idx, total, dark=False):
    fg = RGBColor(0xBF, 0xD0, 0xE8) if dark else COLORS["muted"]
    line_c = RGBColor(0x3A, 0x4A, 0x66) if dark else COLORS["line"]
    _rect(slide, 0.6, 7.14, 12.13, 0.012, line_c)  # 分隔线
    _text(slide, 0.6, 7.20, 9.5, 0.26, [[(META["footer"], 9, False, fg)]])
    _text(slide, 10.9, 7.20, 1.83, 0.26,
          [[("%d / %d" % (idx, total), 9, False, fg)]], align=PP_ALIGN.RIGHT)


def _bg(slide, color=COLORS["bg"]):
    _rect(slide, 0, 0, 13.333, 7.5, color)


def _chip(slide, L, T, W, H, text, fill, text_color, size=12, bold=True):
    _rect(slide, L, T, W, H, fill, round_=True)
    _text(slide, L, T + (H - 0.3) / 2, W, 0.3,
          [[(text, size, bold, text_color)]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


# ─────────────────────────── 页面类型 ───────────────────────────
def add_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.333, 7.5, COLORS["navy"])
    _rect(s, 0, 0, 13.333, 0.14, COLORS["teal"])
    _rect(s, 0.6, 2.0, 1.6, 0.06, COLORS["teal"])
    _text(s, 0.6, 2.25, 12.0, 1.4,
          [[(META["title"], 44, True, COLORS["white"])]])
    _text(s, 0.6, 3.45, 12.0, 0.7,
          [[(META["subtitle"], 22, False, RGBColor(0xBF, 0xD0, 0xE8))]])
    _text(s, 0.6, 4.35, 12.0, 0.5,
          [[("通用底座 ①–⑦  ·  AI 增补层 ⑧–⑪  ·  " + META["version"], 13, False, RGBColor(0x8E, 0xA3, 0xC4))]])
    # 底部能力条
    labels = ["需求·翻译", "规划·落地", "项目·交付", "数据·验证", "战略·方向",
              "协作·影响力", "修养·底色", "AI 认知", "模型评估", "信任安全", "AI 经济"]
    x = 0.6
    y = 6.35
    for lb in labels:
        w = 0.42 + 0.13 * max(0, len(lb) - 2)
        _chip(s, x, y, w, 0.5, lb, RGBColor(0x2A, 0x3D, 0x60),
              RGBColor(0xCF, 0xDD, 0xF0), size=11)
        x += w + 0.18


def add_panorama(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _text(s, 0.6, 0.55, 12.0, 0.9, [[("课程全景：通用底座 + AI 增补层", 30, True, COLORS["ink"])]])
    _text(s, 0.6, 1.4, 12.0, 0.4,
          [[("先看全景建立地图 → ①–⑦ 逐个吃透 → ⑧–⑪ 补 AI 岗位增量 → 附录自评查漏。", 14, False, COLORS["muted"])]])
    _rect(s, 0.6, 1.9, 12.13, 0.03, COLORS["line"])

    def col(L, W, data, accent, tint):
        _rect(s, L, 2.15, W, 4.35, COLORS["card"], line=COLORS["line"])
        _rect(s, L, 2.15, W, 0.7, accent)
        _text(s, L + 0.3, 2.27, W - 0.6, 0.5,
              [[(data["title"], 16, True, COLORS["white"])]], anchor=MSO_ANCHOR.MIDDLE)
        _chip(s, L + W - 1.7, 2.28, 1.4, 0.44, data["tag"],
              tint, accent, size=11)
        y = 3.0
        for num, name, focus in data["items"]:
            _text(s, L + 0.3, y, 0.8, 0.4, [[(num, 16, True, accent)]])
            _text(s, L + 1.15, y, 3.4, 0.4, [[(name, 15, True, COLORS["ink"])]])
            _text(s, L + W - 3.4, y, 3.1, 0.4, [[(focus, 12, False, COLORS["muted"])]],
                  align=PP_ALIGN.RIGHT)
            y += 0.50
        return y

    col(0.6, 5.95, PANORAMA["base"], COLORS["blue"], LAYER_TINT["base"])
    col(6.78, 5.95, PANORAMA["ai"], COLORS["teal"], LAYER_TINT["ai"])
    _text(s, 0.6, 6.7, 12.0, 0.4,
          [[("⑥ 沟通协作是「放大器」让 ①–⑤ 推得动；⑦ 自我修养是「底色」决定天花板。", 12, False, COLORS["muted"])]])
    _footer(s, 2, total)


def add_chapter(prs, ch, idx, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    accent = COLORS["blue"] if ch["layer"] == "base" else COLORS["teal"]
    tint = LAYER_TINT[ch["layer"]]
    # 顶部：序号徽章 + 标题 + 侧重点
    _rect(s, 0.6, 0.55, 0.85, 0.85, accent, round_=True)
    _text(s, 0.6, 0.55, 0.85, 0.85, [[(ch["num"], 26, True, COLORS["white"])]],
          anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    _text(s, 1.65, 0.58, 8.5, 0.8, [[(ch["title"], 28, True, COLORS["ink"])]])
    _chip(s, 10.6, 0.75, 2.13, 0.5, "侧重点 · " + ch["focus"], tint, accent, size=13)
    _rect(s, 0.6, 1.6, 12.13, 0.03, accent)
    # 左列：一句话归纳 + 关键术语
    _text(s, 0.6, 1.85, 5.5, 0.35, [[("一句话归纳", 13, True, accent)]])
    _text(s, 0.6, 2.25, 5.5, 1.35, [[(ch["summary"], 19, True, COLORS["ink"])]])
    _text(s, 0.6, 3.75, 5.5, 0.35, [[("关键术语", 13, True, accent)]])
    term_lines = [[("•  " + t, 13.5, False, COLORS["ink"])] for t in ch["terms"]]
    for i, line in enumerate(term_lines):
        line[0] = (line[0][0], 13.5, False, COLORS["ink"])
    _text(s, 0.6, 4.15, 5.55, 2.4, term_lines)
    # 右列：思维导图（卡片 + 居中图片）
    card_L, card_T, card_W, card_H = 6.35, 1.85, 6.35, 4.75
    _rect(s, card_L, card_T, card_W, card_H, COLORS["card"], line=COLORS["line"])
    img = os.path.join(REPO, ch["img"])
    iw, ih = PILImage.open(img).size
    pad = 0.22
    scale = min((card_W - 2 * pad) / iw, (card_H - 2 * pad) / ih)
    w, h = iw * scale, ih * scale
    s.shapes.add_picture(img, Inches(card_L + (card_W - w) / 2),
                         Inches(card_T + (card_H - h) / 2), Inches(w), Inches(h))
    # 底部 note 条
    note_label, note_text = ch["note"]
    _rect(s, 0.6, 6.64, 12.13, 0.44, tint, round_=True)
    _text(s, 0.85, 6.64, 1.6, 0.44, [[(note_label, 12, True, accent)]],
          anchor=MSO_ANCHOR.MIDDLE)
    _text(s, 2.35, 6.64, 10.3, 0.44, [[(note_text, 11.5, False, COLORS["ink"])]],
          anchor=MSO_ANCHOR.MIDDLE)
    _footer(s, idx, total)


def add_appendix(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _text(s, 0.6, 0.55, 12.0, 0.9, [[("附录：自评 · 快答 · 路线", 30, True, COLORS["ink"])]])
    _rect(s, 0.6, 1.4, 12.13, 0.03, COLORS["line"])
    x = 0.6
    for letter, title, desc in APPENDIX:
        _rect(s, x, 1.75, 3.9, 4.4, COLORS["card"], line=COLORS["line"])
        _rect(s, x, 1.75, 3.9, 0.7, COLORS["navy"])
        _text(s, x + 0.25, 1.87, 0.6, 0.5, [[(letter, 18, True, COLORS["teal"])]],
              anchor=MSO_ANCHOR.MIDDLE)
        _text(s, x + 0.8, 1.87, 3.0, 0.5, [[(title, 14, True, COLORS["white"])]],
              anchor=MSO_ANCHOR.MIDDLE)
        _text(s, x + 0.25, 2.7, 3.4, 3.2, [[(desc, 13, False, COLORS["ink"])]])
        x += 3.9 + 0.21
    _footer(s, 14, total)


def add_end(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.333, 7.5, COLORS["navy"])
    _rect(s, 0, 7.36, 13.333, 0.14, COLORS["teal"])
    _text(s, 0.6, 2.8, 12.0, 1.0, [[("谢谢", 44, True, COLORS["white"])]],
          align=PP_ALIGN.CENTER)
    _text(s, 0.6, 3.9, 12.0, 0.6,
          [[("先把骨架搭起来，往里塞就是了。", 16, False, RGBColor(0xBF, 0xD0, 0xE8))]],
          align=PP_ALIGN.CENTER)
    _footer(s, 15, total, dark=True)


# ─────────────────────────── 主流程 ───────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    total = 2 + len(CHAPTERS) + 2  # 封面 + 全景 + 11章 + 附录 + 结尾 = 15
    add_cover(prs)
    add_panorama(prs, total)
    for i, ch in enumerate(CHAPTERS):
        add_chapter(prs, ch, 3 + i, total)
    add_appendix(prs, total)
    add_end(prs, total)
    prs.save(META["out"])
    print("[OK] generated:", META["out"], "(slides:", len(prs.slides._sldIdLst), ")")


if __name__ == "__main__":
    build()
