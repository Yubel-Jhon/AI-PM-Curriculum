# -*- coding: utf-8 -*-
"""
《AI Agent 行业现状报告》PPT 生成器（21 页 · 券商深度结构 · 嵌图表）
====================================================================
- 结构对齐 5 份券商样本：封面(评级)→摘要→目录→定义→产业链→市场规模→竞争格局→公司对标→
  技术→应用→安全→商业化→趋势→投资建议→风险提示→免责声明。
- 图表：matplotlib 生成的 PNG（ppt/charts/*.png），用「等比缩放」嵌入，杜绝拉伸/溢出。
- 改数据：改 build_charts.py + 本文件各 slide 函数里的文本/表格，渲染逻辑不动。
运行：python ppt/build_charts.py && python ppt/build_agent_report_ppt.py
"""
import os
import struct
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPT_DIR = os.path.dirname(os.path.abspath(__file__))
CHART_DIR = os.path.join(PPT_DIR, "charts")
OUT = os.path.join(PPT_DIR, "AI-Agent行业现状报告.pptx")

COLORS = {
    "navy":  RGBColor(0x1B, 0x2A, 0x4A),
    "blue":  RGBColor(0x2D, 0x6C, 0xDF),
    "teal":  RGBColor(0x0F, 0xA3, 0xA3),
    "ink":   RGBColor(0x1F, 0x29, 0x37),
    "muted": RGBColor(0x6B, 0x72, 0x80),
    "bg":    RGBColor(0xF4, 0xF6, 0xFA),
    "card":  RGBColor(0xFF, 0xFF, 0xFF),
    "line":  RGBColor(0xE5, 0xE7, 0xEB),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "tint_ai": RGBColor(0xE0, 0xF5, 0xF5),
    "tint_blue": RGBColor(0xE8, 0xF0, 0xFE),
    "danger": RGBColor(0xC0, 0x3A, 0x2B),
}
FONT = "Microsoft YaHei"
EMU_W = Inches(13.333)
EMU_H = Inches(7.5)
TOTAL = 21


def _png_size(path):
    with open(path, "rb") as f:
        head = f.read(33)
    if head[:8] != b"\x89PNG\r\n\x1a\n":
        return None
    return struct.unpack(">II", head[16:24])


def _set_font(run, size, bold=False, color=None, name=FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)


def _solid(shape, fill, line=None, line_w=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w or 0.75)
    shape.shadow.inherit = False


def _rect(slide, L, T, W, H, fill, line=None, round_=False, line_w=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(kind, Inches(L), Inches(T), Inches(W), Inches(H))
    if round_:
        try:
            sp.adjustments[0] = 0.08
        except Exception:
            pass
    _solid(sp, fill, line, line_w)
    return sp


def _text(slide, L, T, W, H, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(L), Inches(T), Inches(W), Inches(H))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        for (t, size, bold, color) in para:
            r = p.add_run()
            r.text = t
            _set_font(r, size, bold, color)
    return tb


def _bg(slide, color=COLORS["bg"]):
    _rect(slide, 0, 0, 13.333, 7.5, color)


def _chip(slide, L, T, W, H, text, fill, text_color, size=12, bold=True):
    _rect(slide, L, T, W, H, fill, round_=True)
    _text(slide, L, T, W, H, [[(text, size, bold, text_color)]],
          anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def _footer(slide, idx, total, dark=False):
    fg = RGBColor(0xBF, 0xD0, 0xE8) if dark else COLORS["muted"]
    line_c = RGBColor(0x3A, 0x4A, 0x66) if dark else COLORS["line"]
    _rect(slide, 0.6, 7.14, 12.13, 0.012, line_c)
    _text(slide, 0.6, 7.20, 9.5, 0.26, [[("AI Agent 行业现状报告 · 2026-08", 9, False, fg)]])
    _text(slide, 10.9, 7.20, 1.83, 0.26, [[("%d / %d" % (idx, total), 9, False, fg)]],
          align=PP_ALIGN.RIGHT)


def _header(slide, num, title, tag, accent, idx):
    _bg(slide)
    _rect(slide, 0.6, 0.45, 0.72, 0.72, accent, round_=True)
    _text(slide, 0.6, 0.45, 0.72, 0.72, [[(num, 20, True, COLORS["white"])]],
          anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    _text(slide, 1.5, 0.48, 9.0, 0.7, [[(title, 24, True, COLORS["ink"])]])
    _chip(slide, 10.7, 0.58, 2.0, 0.5, tag, COLORS["tint_ai"], accent, size=12)
    _rect(slide, 0.6, 1.32, 12.13, 0.03, accent)
    _footer(slide, idx, TOTAL)


def _table(slide, L, T, col_widths, headers, rows, header_fill, header_color,
           font_size=11.5, header_size=12, row_h=0.44, header_h=0.5,
           align_right_cols=(), zebra=(COLORS["card"], RGBColor(0xF0, 0xF3, 0xF8))):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    W = sum(col_widths)
    gfx = slide.shapes.add_table(n_rows, n_cols, Inches(L), Inches(T),
                                 Inches(W), Inches(header_h + row_h * len(rows))).table
    for i, cw in enumerate(col_widths):
        gfx.columns[i].width = Inches(cw)
    gfx.rows[0].height = Inches(header_h)
    for i in range(len(rows)):
        gfx.rows[i + 1].height = Inches(row_h)

    def _cell(cell, text, size, bold, color, fill, align):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        _set_font(r, size, bold, color)

    for c, h in enumerate(headers):
        _cell(gfx.cell(0, c), h, header_size, True, header_color, header_fill, PP_ALIGN.LEFT)
    for ri, row in enumerate(rows):
        fill = zebra[ri % 2]
        for ci, val in enumerate(row):
            align = PP_ALIGN.RIGHT if ci in align_right_cols else PP_ALIGN.LEFT
            _cell(gfx.cell(ri + 1, ci), val, font_size, False, COLORS["ink"], fill, align)
    return gfx


def _pic(slide, name, L, T, W, H):
    """等比缩放居中嵌入图表，杜绝拉伸/溢出。"""
    path = os.path.join(CHART_DIR, name + ".png")
    if not os.path.exists(path):
        return None
    w, h = _png_size(path)
    ratio = w / h
    pw, ph = W, W / ratio
    if ph > H:
        ph, pw = H, H * ratio
    left = L + (W - pw) / 2
    top = T + (H - ph) / 2
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(pw), Inches(ph))
    return (left, top, pw, ph)


# ─────────────────────────── 页面 ───────────────────────────
def add_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.333, 7.5, COLORS["navy"])
    _rect(s, 0, 0, 13.333, 0.14, COLORS["teal"])
    _chip(s, 0.6, 0.7, 2.4, 0.5, "行业深度研究报告", RGBColor(0x2A, 0x3D, 0x60),
          RGBColor(0xCF, 0xDD, 0xF0), size=12)
    _chip(s, 3.15, 0.7, 1.6, 0.5, "首次覆盖", RGBColor(0x2A, 0x3D, 0x60),
          RGBColor(0x8E, 0xF2, 0xE8), size=12)
    _rect(s, 0.6, 2.0, 1.6, 0.06, COLORS["teal"])
    _text(s, 0.6, 2.2, 12.0, 1.3, [[("AI Agent 行业现状报告", 42, True, COLORS["white"])]])
    _text(s, 0.6, 3.5, 12.0, 0.9,
          [[("从「能不能做」到「算不算得过账」—— Agent 进入「做得好不好、算不算得过来、出了事谁负责」的验证期",
             19, False, RGBColor(0xBF, 0xD0, 0xE8))]])
    _chip(s, 0.6, 4.6, 2.2, 0.55, "行业评级：看好", RGBColor(0x2A, 0x3D, 0x60),
          RGBColor(0xCF, 0xDD, 0xF0), size=14)
    _text(s, 0.6, 5.4, 12.0, 0.4,
          [[("整理：AI-PM-Curriculum 研究组　|　报告日期：2026-08-15", 12.5, False, RGBColor(0x8E, 0xA3, 0xC4))]])
    chips = ["摘要", "定义", "产业链", "市场规模", "竞争格局", "公司对标", "技术", "应用", "安全", "商业化", "趋势", "风险"]
    x = 0.6
    y = 6.35
    for c in chips:
        w = 0.5 + 0.155 * max(0, len(c) - 2)
        if x + w > 12.7:
            break
        _chip(s, x, y, w, 0.46, c, RGBColor(0x2A, 0x3D, 0x60), RGBColor(0xCF, 0xDD, 0xF0), size=10.5)
        x += w + 0.14


def add_summary(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "0", "摘要 · 核心观点", "先看结论", COLORS["teal"], idx)
    _rect(s, 0.6, 1.5, 12.13, 0.9, COLORS["navy"], round_=True)
    _text(s, 0.85, 1.56, 11.7, 0.8,
          [[("核心结论：", 13, True, COLORS["teal"]),
            ("Agent 已从「能不能做」进入「做得好不好、算不算得过来、出了事谁负责」的验证期；"
             "技术撑得起单点场景，但可靠性、单位经济、责任归属三道坎，让它仍卡在「单点工具」。",
             13, False, COLORS["white"])]])
    _text(s, 0.6, 2.55, 6.3, 0.4, [[("五条判断", 14, True, COLORS["ink"])]])
    judgments = [
        "瓶颈从「模型不够聪明」转向「不够可靠、可控、便宜」。",
        "市场「框架热、应用冷」：收得上钱的仍是代码 / 客服 / 办公。",
        "竞争从「模型之争」转向「工具 + 数据闭环之争」。",
        "商业化「按调用 / 按结果付费」为主，多数单位经济未打正。",
        "安全合规是「发布门槛」，责任不清是 B 端采购最大隐性阻力。",
    ]
    _text(s, 0.6, 2.95, 6.3, 2.6, [[("▸ " + j, 12, False, COLORS["ink"])] for j in judgments], spacing=1.22)
    _text(s, 7.0, 2.55, 5.7, 0.4, [[("关键数据快照", 14, True, COLORS["ink"])]])
    _table(s, 7.0, 2.95, [2.35, 3.35],
           ["指标", "量级判断"],
           [["狭义规模", "2024 52.6亿 → 2030 526亿美元"],
            ["广义支出", "Gartner 2026 2019亿 → 2029 7527亿"],
            ["编码能力", "SWE-bench 80.9%（Claude Opus 4.5）"],
            ["多步成功率", "单步 90% × 20 步 → 约 12%"],
            ["头部产品 ARR", "Cursor $4B · Claude Code $2.5B"]],
           COLORS["teal"], COLORS["white"], font_size=10, row_h=0.4)
    _rect(s, 0.6, 5.75, 6.3, 1.15, COLORS["tint_ai"], round_=True)
    _text(s, 0.8, 5.83, 5.95, 0.4, [[("投资建议（简版）", 12, True, COLORS["teal"])]])
    _text(s, 0.8, 6.2, 5.95, 0.65,
          [[("沿「代码→客服→办公→金融/法律」成熟度顺序，优先垂直 Agent 与有数据飞轮的场景。",
             11, False, COLORS["ink"])]])
    _rect(s, 7.0, 5.75, 5.7, 1.15, COLORS["tint_blue"], round_=True)
    _text(s, 7.2, 5.83, 5.3, 0.4, [[("风险提示（简版）", 12, True, COLORS["danger"])]])
    _text(s, 7.2, 6.2, 5.3, 0.65,
          [[("可靠性不及预期；单位经济打不正；提示注入 / 供应链投毒；监管趋严；同质化。",
             11, False, COLORS["ink"])]])


def add_toc(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "·", "目录", "结构总览", COLORS["blue"], idx)
    items = [
        ("一", "行业复盘与定义", "Agent 是什么、从哪来"),
        ("二", "产业链分析", "上游算力/模型 → 中游框架 → 下游应用"),
        ("三", "市场规模与测算", "口径拆解 + 逐年增长 + 区域/细分"),
        ("四", "竞争格局", "份额、集中度、价值链、竞争矩阵"),
        ("五", "关键公司对标", "8–12 家头部公司 + ARR 对比"),
        ("六", "技术现状", "透镜⑧：能不能做"),
        ("七", "应用落地", "透镜⑨：做得好不好"),
        ("八", "信任安全合规", "透镜⑩：守住底线"),
        ("九", "商业化经济性", "透镜⑪：值不值赚不赚"),
        ("十", "趋势判断", "12–24 个月"),
        ("十一", "投资建议", "评级 + 按环节列标的"),
        ("十二", "风险提示", "条目式"),
    ]
    for i, (num, title, desc) in enumerate(items):
        col = i // 6
        row = i % 6
        L = 0.6 + col * 6.2
        T = 1.6 + row * 0.86
        _rect(s, L, T, 0.55, 0.55, COLORS["blue"], round_=True)
        _text(s, L, T, 0.55, 0.55, [[(num, 16, True, COLORS["white"])]],
              anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        _text(s, L + 0.72, T + 0.02, 5.15, 0.4, [[(title, 15, True, COLORS["ink"])]])
        _text(s, L + 0.72, T + 0.4, 5.15, 0.35, [[(desc, 10.5, False, COLORS["muted"])]])


def add_definition(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "1", "一、行业复盘与定义", "从哪来 · 是什么", COLORS["blue"], idx)
    _rect(s, 0.6, 1.5, 6.0, 1.05, COLORS["tint_blue"], round_=True)
    _text(s, 0.8, 1.57, 5.6, 0.4, [[("定义", 13, True, COLORS["blue"])]])
    _text(s, 0.8, 1.94, 5.6, 0.6,
          [[("规划 + 工具调用 + 记忆 + 自主执行——拆解任务、调用工具、循环逼近目标，不再只是回答问题。",
             12, False, COLORS["ink"])]])
    _text(s, 0.6, 2.72, 6.0, 0.4, [[("与相邻形态：自主程度连续谱", 14, True, COLORS["ink"])]])
    _table(s, 0.6, 3.12, [1.0, 1.4, 0.7, 1.0],
           ["形态", "特征", "自主", "关系"],
           [["Chatbot", "只输出文本", "无", "底座"],
            ["Copilot", "人主导", "低", "前置"],
            ["RPA/工作流", "规则固定", "中", "互补"],
            ["AI Agent", "规划+工具+循环", "高", "收敛点"]],
           COLORS["blue"], COLORS["white"], font_size=10.5, row_h=0.44)
    _text(s, 7.0, 1.5, 5.7, 0.4, [[("发展历程（2022.11 → 2026.04）", 14, True, COLORS["ink"])]])
    _table(s, 7.0, 1.92, [1.25, 4.45],
           ["时间", "事件"],
           [["2022.11", "ChatGPT 发布"],
            ["2023.06", "Function Calling 工具调用 API 化"],
            ["2024.09/11", "o1 推理 / MCP 开源"],
            ["2025.01/04", "Operator / Google A2A"],
            ["2025.12", "MCP 捐 Linux 基金会"],
            ["2026.04", "Gartner Hype Cycle 首发布"]],
           COLORS["blue"], COLORS["white"], font_size=10.5, row_h=0.55)
    _rect(s, 0.6, 5.75, 12.13, 0.9, COLORS["tint_blue"], round_=True)
    _text(s, 0.85, 5.85, 11.7, 0.7,
          [[("复盘：", 12.5, True, COLORS["ink"]),
            ("Copilot → Agent → Multi-Agent；2025H2–2026H1 从「跑马圈地」进入「收敛」——标准定型、框架 1.0 化、编码 Agent 跑通商业模式、监管入场。",
             12.5, False, COLORS["ink"])]], spacing=1.2)


def add_industry_chain(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "2", "二、产业链分析", "上中下游", COLORS["blue"], idx)
    cols = [
        ("上游 · 算力/模型", COLORS["teal"],
         [("算力", "Nvidia 份额 80–90%，数据中心收入约 1940 亿美元/年；推理取代训练成增长主引擎"),
          ("模型", "OpenAI $40B · Anthropic $30B · Google · DeepSeek · 智谱 · 通义 · 豆包"),
          ("数据", "行业 know-how + 私域数据，价值沉淀在垂直应用层")]),
        ("中游 · 框架/平台", COLORS["blue"],
         [("编排", "LangChain（月下载 9000 万）· CrewAI · Microsoft Agent Framework"),
          ("平台", "Copilot Studio（80%+ F500）· Dify · Coze · 百度千帆 · 阿里百炼"),
          ("特征", "薄利润、强网络效应；开源免费是默认，护城河最弱")]),
        ("下游 · 垂直应用", COLORS["teal"],
         [("代码", "市场约 128 亿美元，CR3 占 70–80%；最成熟"),
          ("客服", "Agentforce $1.2B · Sierra · Fin；按结果付费跑通"),
          ("办公/金融/法律", "Copilot 3000 万席位 · Harvey 等，成熟度递减")]),
    ]
    x = 0.6
    for name, accent, items in cols:
        _rect(s, x, 1.5, 3.95, 0.5, accent, round_=True)
        _text(s, x, 1.5, 3.95, 0.5, [[(name, 13.5, True, COLORS["white"])]],
              anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        y = 2.12
        for label, desc in items:
            _rect(s, x, y, 3.95, 1.05, COLORS["card"], line=COLORS["line"], round_=True)
            _text(s, x + 0.16, y + 0.08, 3.65, 0.35, [[(label, 11.5, True, accent)]])
            _text(s, x + 0.16, y + 0.42, 3.65, 0.6, [[(desc, 9.5, False, COLORS["ink"])]])
            y += 1.17
        x += 4.09
    _rect(s, 0.6, 5.85, 12.13, 0.95, COLORS["tint_ai"], round_=True)
    _text(s, 0.85, 5.95, 11.7, 0.75,
          [[("传导逻辑：", 12.5, True, COLORS["teal"]),
            ("开源模型 + 混合路由使企业 token 成本 1 年降 67%（$18.40→$6.07/M）→ 上游降价打开下游毛利空间 → 下游按结果计价又放大上游 token 消耗，正反馈闭环。",
             12, False, COLORS["ink"])]], spacing=1.15)


def add_market_scale(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "3", "三、市场规模与测算（1/2）", "增长曲线", COLORS["teal"], idx)
    _rect(s, 0.6, 1.45, 12.13, 0.55, COLORS["tint_ai"], round_=True)
    _text(s, 0.8, 1.45, 11.7, 0.55,
          [[("口径拆解：", 11.5, True, COLORS["teal"]),
            ("狭义（软件/平台）MarketsandMarkets/GVR/BCC ≈ 2025 78–80亿、2030 483–526亿；广义（Gartner）2029 7527亿，差一个数量级。",
             11.5, False, COLORS["ink"])]], anchor=MSO_ANCHOR.MIDDLE)
    _pic(s, "fig01_market", 0.6, 2.15, 6.0, 4.2)
    _pic(s, "fig02_gartner", 6.75, 2.15, 6.0, 4.2)
    _text(s, 0.6, 6.45, 6.0, 0.5,
          [[("狭义口径三家收敛：CAGR 43%–46%。", 11, True, COLORS["ink"])]])
    _text(s, 6.75, 6.45, 6.0, 0.5,
          [[("Gartner：2027 年 Agentic 支出超越 chatbot；40% 项目或 2027 底前被搁置。", 11, True, COLORS["ink"])]])


def add_market_segment(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "3", "三、市场规模与测算（2/2）", "区域/细分", COLORS["teal"], idx)
    _pic(s, "fig03_region", 0.6, 1.6, 6.0, 4.3)
    _pic(s, "fig04_vertical", 6.75, 1.6, 6.0, 4.3)
    _text(s, 0.6, 6.0, 6.0, 0.8,
          [[("区域：北美最大（44%），亚太增速最快。", 11.5, True, COLORS["ink"])]])
    _text(s, 6.75, 6.0, 6.0, 0.8,
          [[("垂直 Agent（62.7%）显著快于通用（44.9%）；5 个垂直吃掉 520 亿市场的 80%。",
             11.5, True, COLORS["ink"])]])


def add_competition_share(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "4", "四、竞争格局（1/2）", "份额/集中度", COLORS["blue"], idx)
    _pic(s, "fig05_platform", 0.6, 1.6, 6.0, 4.3)
    _pic(s, "fig06_model_share", 6.75, 1.6, 6.0, 4.3)
    _text(s, 0.6, 6.0, 6.0, 0.9,
          [[("企业平台部署 CR3 73%：Microsoft 31% + Salesforce 24% + Anthropic 18%。",
             11.5, True, COLORS["ink"])],
           [("按收入计集中度反而更低（CR5 略超 50%），74% 企业担忧锁定。", 10.5, False, COLORS["muted"])]], spacing=1.15)
    _text(s, 6.75, 6.0, 6.0, 0.9,
          [[("模型层企业工作负载：Anthropic 40% 反超 OpenAI 27%。", 11.5, True, COLORS["ink"])],
           [("编码 Agent CR3 占 70–80%；火山 MaaS 国内份额 49.5%。", 10.5, False, COLORS["muted"])]], spacing=1.15)


def add_competition_matrix(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "4", "四、竞争格局（2/2）", "价值链/矩阵", COLORS["blue"], idx)
    _pic(s, "fig07_valuechain", 0.6, 1.6, 5.7, 4.3)
    _text(s, 0.6, 6.0, 5.7, 0.8,
          [[("算力仍是最大现金池（英伟达单季 752 亿 > 模型厂+应用厂之和）——「卖铲人吃肉」。",
             11, True, COLORS["ink"])]])
    _table(s, 6.5, 1.6, [1.05, 1.70, 1.70, 1.70],
           ["维度", "上游模型", "中游平台", "下游应用"],
           [["竞争烈度", "高", "中高", "中"],
            ["护城河", "模型+算力+数据", "生态锁定(弱)", "know-how+数据"],
            ["赢家通吃", "强", "弱", "弱-中"],
            ["代表", "OpenAI等", "LangChain等", "Agentforce等"]],
           COLORS["blue"], COLORS["white"], font_size=10, row_h=0.5)
    _text(s, 6.5, 6.0, 6.2, 0.8,
          [[("护城河沿产业链向下游「行业 know-how + 私域数据」迁移；中游最弱，随时被上下游夹击。",
             11, True, COLORS["ink"])]])


def add_companies_table(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "5", "五、关键公司对标（1/2）", "对标表", COLORS["blue"], idx)
    _table(s, 0.6, 1.5, [1.7, 2.6, 2.9, 2.35, 2.6],
           ["公司", "核心产品", "商业化(2026)", "用户/估值", "护城河"],
           [["OpenAI", "ChatGPT/Codex", "年化 $40B", "周活9亿·$852B", "模型+用户+全栈"],
            ["Anthropic", "Claude/Claude Code", "ARR $30B", "$380B", "企业信任+编码"],
            ["Microsoft", "Copilot/Studio", "AI运行率 $37B", "席位3000万", "Office分发+Azure"],
            ["Salesforce", "Agentforce", "ARR $1.2B(+205%)", "2.9万订单", "CRM数据+交叉销售"],
            ["Cursor", "Cursor", "ARR ~$4B", "$60B收购", "编码IDE市占"],
            ["Sierra/Fin", "客服Agent", "$0.2B/$0.1B", "$15.8B/$3.6B", "按结果付费"]],
           COLORS["blue"], COLORS["white"], font_size=10, row_h=0.6)
    _text(s, 0.6, 5.95, 12.13, 0.8,
          [[("国内：智谱 MaaS 17亿RMB · 豆包 ARR ~$4B · 通义百炼 80亿RMB，均未跑出「Agent 专属收入」大数。",
             11.5, True, COLORS["ink"])]])


def add_companies_chart(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "5", "五、关键公司对标（2/2）", "ARR 对比", COLORS["blue"], idx)
    _pic(s, "fig08_arr", 1.4, 1.6, 10.5, 4.6)
    _text(s, 0.6, 6.3, 12.13, 0.7,
          [[("结论：商业化分三层——模型+平台赚大钱（OpenAI $40B/Microsoft $37B/Anthropic $30B）；"
             "垂直 Agent 高增速小基数（Cursor $4B/Agentforce $1.2B）；按结果付费客服最小但最快（Fin +350%）。",
             11.5, True, COLORS["ink"])]])


def add_tech_bench(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "6", "六、技术现状（1/3）", "能不能做 · ⑧", COLORS["teal"], idx)
    _pic(s, "fig09_swebench", 0.6, 1.6, 6.0, 4.3)
    _pic(s, "fig10_osworld", 6.75, 1.6, 6.0, 4.3)
    _text(s, 0.6, 6.0, 6.0, 0.8,
          [[("SWE-bench：Claude Opus 4.5 达 80.9%，编码是提升最快、可验证的领域。",
             11, True, COLORS["ink"])]])
    _text(s, 6.75, 6.0, 6.0, 0.8,
          [[("OSWorld：Simular S2 72.6% 首次超人类（72.36%）；GUI 操作仍在爬坡。",
             11, True, COLORS["ink"])]])


def add_tech_metr_cost(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "6", "六、技术现状（2/3）", "长程/成本 · ⑧", COLORS["teal"], idx)
    _pic(s, "fig11_metr", 0.6, 1.6, 6.0, 4.3)
    _pic(s, "fig12_cost", 6.75, 1.6, 6.0, 4.3)
    _text(s, 0.6, 6.0, 6.0, 0.8,
          [[("METR：50% 时间线约每 3–4 个月翻倍（Opus 4.5 达 289 分钟）；要求 80% 成功率时骤降至 15 分钟。",
             11, True, COLORS["ink"])]])
    _text(s, 6.75, 6.0, 6.0, 0.8,
          [[("成本：旗舰输出价 17 个月降 6 倍（$60→$10），小模型侧降 100 倍。",
             11, True, COLORS["ink"])]])


def add_tech_decay(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "6", "六、技术现状（3/3）", "复利误差 · ⑧", COLORS["teal"], idx)
    _pic(s, "fig13_decay", 0.6, 1.6, 6.4, 4.3)
    _text(s, 7.3, 1.6, 5.4, 0.4, [[("技术瓶颈", 14, True, COLORS["ink"])]])
    bottlenecks = [
        "幻觉：通用 22%–94%，执行类代价被急剧放大",
        "多步成功率：单步 90% × 20 步 → 仅约 12%",
        "长程任务：真实成功率多数 < 20%（METR）",
        "延迟成本：多步 p95 达 58–90 秒",
    ]
    _text(s, 7.3, 2.05, 5.4, 2.6, [[("▸ " + b, 12, False, COLORS["ink"])] for b in bottlenecks], spacing=1.4)
    _rect(s, 7.3, 4.9, 5.4, 1.5, COLORS["tint_ai"], round_=True)
    _text(s, 7.5, 5.0, 5.0, 1.3,
          [[("判断：", 12.5, True, COLORS["teal"]),
            ("单点、封闭、可验证场景「能」；基准快速上升，真实环境端到端慢一个量级。",
             12, False, COLORS["ink"])]], spacing=1.2)


def add_application(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "7", "七、应用落地", "做得好不好 · ⑨", COLORS["teal"], idx)
    _table(s, 0.6, 1.5, [0.6, 1.7, 1.8, 1.5],
           ["场景", "落地形态", "为何跑通", "关键数字"],
           [["代码", "Claude Code/Cursor", "可编译可测试", "SWE-bench + ARR"],
            ["客服", "Agentforce/Fin", "封闭+转人工兜底", "解决率63–90%"],
            ["办公", "M365 Copilot", "流程相对确定", "席位3000万+"],
            ["营销", "万相台/AgenticOS", "ROI可归因", "ROI+20%"],
            ["金融", "Bloomberg ASKB", "数据密集+人机协同", "尽调周级→小时"]],
           COLORS["teal"], COLORS["white"], font_size=10, row_h=0.5)
    _pic(s, "fig14_penetration", 6.6, 1.5, 6.0, 3.4)
    _rect(s, 0.6, 5.5, 6.0, 1.3, COLORS["card"], line=COLORS["line"], round_=True)
    _text(s, 0.8, 5.58, 5.6, 0.4, [[("仍未跑通", 12.5, True, COLORS["danger"])]])
    _text(s, 0.8, 5.95, 5.6, 0.8,
          [[("长程自主任务（METR「0% 无需人工清理可用」）；高风险执行；长记忆/跨系统状态。",
             10.5, False, COLORS["ink"])]])
    _text(s, 6.6, 5.5, 6.0, 1.3,
          [[("分水岭不是「有没有 AI 价值」，而是「能否建评估集 + 在线指标 + 数据闭环」——评测能力决定落地顺序。",
             11, True, COLORS["ink"])]], anchor=MSO_ANCHOR.MIDDLE)


def add_trust(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "8", "八、信任、安全与合规", "守住底线 · ⑩", COLORS["teal"], idx)
    risks = [
        ("提示注入（头号）", "OWASP LLM01 连续两版第一；受控测试凭据外泄 96%"),
        ("越权/误操作", "Agent 继承用户完整权限，缺最小权限控制"),
        ("MCP 供应链投毒", "工具描述中毒平均攻击成功率 36.5%"),
    ]
    x = 0.6
    for name, desc in risks:
        _rect(s, x, 1.5, 3.95, 1.0, COLORS["card"], line=COLORS["line"], round_=True)
        _text(s, x + 0.2, 1.6, 3.55, 0.4, [[(name, 12.5, True, COLORS["danger"])]])
        _text(s, x + 0.2, 1.98, 3.55, 0.5, [[(desc, 10, False, COLORS["ink"])]])
        x += 4.09
    _text(s, 0.6, 2.7, 6.0, 0.4, [[("责任归属", 14, True, COLORS["ink"])]])
    _text(s, 0.6, 3.1, 6.2, 1.5,
          [[("AI 无法律人格，责任落在部署方/开发方（Air Canada 被判兑现聊天机器人承诺）。",
             12, False, COLORS["ink"])],
           [("「法律不确定、部署方兜底、传统保险退场、专属 AI 险萌芽」。", 12, True, COLORS["ink"])]],
          spacing=1.25)
    _text(s, 7.0, 2.7, 5.7, 0.4, [[("监管三地分轨", 14, True, COLORS["ink"])]])
    _text(s, 7.0, 3.1, 5.7, 1.5,
          [[("中国：双备案 + 内容标识 + 数据出境（已处置违规智能体 3500+ 款）。", 11, False, COLORS["ink"])],
           [("欧盟：EU AI Act「可审计 + 可兜底 + 可问责」。", 11, False, COLORS["ink"])],
           [("美国：联邦轻监管、州法反推、DOJ 刑事执法优先。", 11, False, COLORS["ink"])]],
          spacing=1.25)
    _rect(s, 0.6, 5.0, 12.13, 1.5, COLORS["tint_ai"], round_=True)
    _text(s, 0.85, 5.15, 11.7, 0.4, [[("发布门槛：安全 + 兜底 + 指标 + 监控，不过关不发", 13, True, COLORS["teal"])]])
    _text(s, 0.85, 5.6, 11.7, 0.8,
          [[("业界正从「模型层对齐」转向「环境层强制 + 最小权限 + 人工确认关键动作 + 可审计日志」。",
             12.5, False, COLORS["ink"])],
           [("Agent 产品经理第一职责：先想清楚怎么不闯祸，再谈功能多强。", 12.5, True, COLORS["ink"])]], spacing=1.2)


def add_business(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "9", "九、商业化与经济性", "值不值赚不赚 · ⑪", COLORS["blue"], idx)
    _rect(s, 0.6, 1.5, 12.13, 0.68, COLORS["tint_blue"], round_=True)
    _text(s, 0.85, 1.56, 11.7, 0.58,
          [[("一次任务 = 规划 + 多轮工具调用 + 反思 + 重试，Token 放大 5–30x（编码极端 1000x）；主因是重试与静默失败，不是模型贵。",
             12.5, False, COLORS["ink"])]], anchor=MSO_ANCHOR.MIDDLE)
    _table(s, 0.6, 2.4, [1.0, 2.8, 2.0],
           ["环节", "现状", "判断"],
           [["推理成本", "同能力 2 年降 80–90%", "被多步放大抵消"],
            ["定价建模", "向按结果付费迁移", "把风险转给供应商"],
            ["付费意愿", "B 端愿付，C 端 71% 要免费", "分化大"],
            ["数据飞轮", "代码/客服有，多数没有", "真正的护城河"]],
           COLORS["blue"], COLORS["white"], font_size=11, row_h=0.46)
    _rect(s, 7.0, 2.4, 5.7, 3.0, COLORS["navy"], round_=True)
    _text(s, 7.25, 2.55, 5.2, 0.4, [[("核心矛盾", 13, True, COLORS["teal"])]])
    _text(s, 7.25, 3.0, 5.2, 2.3,
          [[("价值越高（越自主）→ 单次成本越高 → 越难定价 → 越难算清账。", 13, True, COLORS["white"])],
           [("单位经济只在「高频、可量化、容错低」场景打正（代码/客服）。", 12, False, RGBColor(0xBF, 0xD0, 0xE8))],
           [("→ 选场景即选商业模式，第一性动作是「算账」而非「堆能力」。", 13, True, COLORS["teal"])]],
          spacing=1.3)
    _rect(s, 0.6, 5.85, 12.13, 0.8, COLORS["tint_blue"], round_=True)
    _text(s, 0.85, 5.9, 11.7, 0.7,
          [[("按结果付费：HubSpot $0.50/解决 · Intercom Fin $0.99/解决 · Sierra $1.50/解决——护城河来自数据飞轮，不是算法。",
             12, True, COLORS["ink"])]], anchor=MSO_ANCHOR.MIDDLE)


def add_trends(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "10", "十、趋势判断", "12–24 个月", COLORS["teal"], idx)
    _rect(s, 0.6, 1.5, 6.0, 0.55, COLORS["teal"], round_=True)
    _text(s, 0.6, 1.5, 6.0, 0.55, [[("高确定性趋势", 14, True, COLORS["white"])]],
          anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    certain = [
        "推理成本持续下降（12 个月再降 40%+），但多步吃掉红利",
        "工具调用标准化：MCP / A2A 成事实标准",
        "Agent 嵌入垂直工作流（垂直 CAGR 62.7% > 整体 44.9%）",
        "安全合规前置：可审计、可兜底、分层发布成标配",
        "「慢思考 + 反思」成默认架构，按难度分配算力",
    ]
    _text(s, 0.7, 2.2, 5.8, 4.3, [[("▸ " + t, 11.5, False, COLORS["ink"])] for t in certain], spacing=1.5)
    _rect(s, 6.85, 1.5, 5.9, 0.55, COLORS["danger"], round_=True)
    _text(s, 6.85, 1.5, 5.9, 0.55, [[("高分歧 / 不确定点", 14, True, COLORS["white"])]],
          anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    uncertain = [
        "多智能体是否必要（默认单 Agent，必要时才多 Agent）",
        "是否出现「Agent 平台型超级应用」，入口归属未定",
        "通用 Agent 何时可托付（乐观 2026 底–2027，谨慎半自主）",
        "按结果付费能否主流（混合模式更现实，纯结果短期难）",
    ]
    _text(s, 6.95, 2.2, 5.7, 4.3, [[("▸ " + t, 11.5, False, COLORS["ink"])] for t in uncertain], spacing=1.5)


def add_investment(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "11", "十一、投资建议", "评级 + 标的", COLORS["blue"], idx)
    _rect(s, 0.6, 1.5, 12.13, 0.85, COLORS["navy"], round_=True)
    _text(s, 0.85, 1.56, 1.6, 0.4, [[("行业评级", 13, True, COLORS["teal"])]])
    _text(s, 0.85, 1.92, 11.6, 0.4,
          [[("看好（首次覆盖）：技术单点可用 + 代码/客服商业模式已跑通 + 成本持续下探，结构性机会在「垂直 Agent + 数据飞轮」。",
             13, False, COLORS["white"])]])
    _table(s, 0.6, 2.6, [2.3, 3.2, 3.3],
           ["环节", "关注", "理由"],
           [["模型/算力层", "OpenAI · Anthropic · Nvidia", "算力为最大现金池"],
            ["垂直 Agent 层", "Cursor · Agentforce · Sierra · Fin", "代码/客服已跑通商业化"],
            ["国内", "智谱 · 豆包 · 通义", "MaaS 层，Agent 专属收入待验证"]],
           COLORS["blue"], COLORS["white"], font_size=11, row_h=0.5)
    advice = [
        ("1", "用「评测」倒逼「落地」", "先定义评估集与在线指标，评不出的场景先不做（⑨）。"),
        ("2", "把「单位经济」当第一过滤器", "优先高频、可量化、容错低、有数据飞轮的场景（⑪）。"),
        ("3", "把安全与责任当「发布门槛」", "内建权限边界、兜底与可审计，红队 + 分层发布（⑩）。"),
    ]
    y = 4.9
    for num, title, desc in advice:
        _rect(s, 0.6, y, 12.13, 0.52, COLORS["card"], line=COLORS["line"], round_=True)
        _rect(s, 0.6, y, 0.6, 0.52, COLORS["blue"], round_=True)
        _text(s, 0.6, y, 0.6, 0.52, [[(num, 17, True, COLORS["white"])]],
              anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        _text(s, 1.4, y, 11.2, 0.52,
              [[(title, 13, True, COLORS["ink"]), ("　" + desc, 12, False, COLORS["ink"])]],
              anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT)
        y += 0.60


def add_risk(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "12", "十二、风险提示", "条目式", COLORS["blue"], idx)
    risks = [
        ("多步可靠性提升不及预期", "单步衰减是数学约束，若长程成功率无法突破，Agent 将长期停留在「单点工具」。"),
        ("单位经济长期打不正", "多步放大推理成本，若按结果付费覆盖不了重试成本，商业化将停滞。"),
        ("提示注入 / 供应链投毒", "头号威胁无纯模型解，重大事故将触发监管收紧与采购观望。"),
        ("监管趋严抬高门槛", "EU AI Act / 中国双备案若加码，抬高发布成本、拖慢落地。"),
        ("同质化竞争价值难沉淀", "中游框架供给过剩、开源免费，利润被上下游两头挤压。"),
        ("市场空间测算偏差", "规模测算基于第三方口径与 CAGR 假设，实际可能与测算存在偏差。"),
    ]
    y = 1.5
    for name, desc in risks:
        _rect(s, 0.6, y, 12.13, 0.82, COLORS["card"], line=COLORS["line"], round_=True)
        _rect(s, 0.6, y, 0.07, 0.82, COLORS["danger"])
        _text(s, 0.85, y + 0.08, 3.4, 0.4, [[(name, 13, True, COLORS["ink"])]])
        _text(s, 0.85, y + 0.46, 11.7, 0.35, [[(desc, 11, False, COLORS["muted"])]])
        y += 0.92


def add_disclaimer(prs, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.333, 7.5, COLORS["navy"])
    _rect(s, 0, 7.36, 13.333, 0.14, COLORS["teal"])
    _text(s, 0.6, 0.7, 12.0, 0.8, [[("免责声明 · 评级说明 · 分析师声明", 26, True, COLORS["white"])]])
    _text(s, 0.6, 1.7, 12.1, 3.2,
          [[("免责声明：", 13, True, COLORS["teal"]),
            ("本报告基于公开信息与第三方机构（MarketsandMarkets、GVR、BCC、Gartner、Omdia、"
             "Papers with Code、METR、Menlo Ventures 等）数据整理，仅供学习与研究参考，不构成任何投资建议。"
             "文中标注 [估计]/[推算] 的数据为基于公开 CAGR 或跨口径推断的估算值，可能与实际存在偏差；"
             "数据截至 2026-08，如有最新一手来源应回填替换。", 12.5, False, RGBColor(0xE8, 0xEE, 0xF7))],
           [("评级说明：", 13, True, COLORS["teal"]),
            ("买入（相对指数 +15% 以上）/ 增持（5%–15%）/ 持有（-5%–5%）/ 卖出（-5% 以下）；"
             "行业评级：强于大市 / 中性 / 弱于大市。本报告为行业层「看好」，不对单一标的给出个股评级。",
             12.5, False, RGBColor(0xE8, 0xEE, 0xF7))],
           [("分析师声明：", 13, True, COLORS["teal"]),
            ("本报告作者不持有文中提及任何公司的证券，与所述公司无利益关联。未经许可，不得转载或用于商业用途。",
             12.5, False, RGBColor(0xE8, 0xEE, 0xF7))]],
           spacing=1.5)
    _text(s, 0.6, 5.9, 12.0, 0.9, [[("谢谢", 40, True, COLORS["white"])]],
          align=PP_ALIGN.CENTER)
    _footer(s, idx, TOTAL, dark=True)


# ─────────────────────────── 主流程 ───────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    add_cover(prs)
    add_summary(prs, 2)
    add_toc(prs, 3)
    add_definition(prs, 4)
    add_industry_chain(prs, 5)
    add_market_scale(prs, 6)
    add_market_segment(prs, 7)
    add_competition_share(prs, 8)
    add_competition_matrix(prs, 9)
    add_companies_table(prs, 10)
    add_companies_chart(prs, 11)
    add_tech_bench(prs, 12)
    add_tech_metr_cost(prs, 13)
    add_tech_decay(prs, 14)
    add_application(prs, 15)
    add_trust(prs, 16)
    add_business(prs, 17)
    add_trends(prs, 18)
    add_investment(prs, 19)
    add_risk(prs, 20)
    add_disclaimer(prs, 21)
    prs.save(OUT)
    print("[OK] generated:", OUT, "(slides:", len(prs.slides._sldIdLst), ")")


if __name__ == "__main__":
    build()
